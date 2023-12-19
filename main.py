import sys
import io
import xlsxwriter
from docx import Document
from pptx import Presentation
from docx.shared import Inches, Mm
from random import randint, choice
from datetime import datetime, time, timedelta

from PyQt5 import uic  # Импортируем uic
from PyQt5.QtWidgets import QApplication, QMainWindow, QWidget, QInputDialog, QFileDialog, QPushButton
import PyQt5

from modules import *


class MyCinemasSystem(QMainWindow):
    def __init__(self):
        super().__init__()
        uic.loadUi("main.ui", self)

        self.cinemas = {}
        self.sessions = []

        for s in self.sessions:
            for i in range(len(s['hall'])):
                for j in range(len(s['hall'][0])):
                    s['hall'][i][j] = randint(0, 1)
        self.addCinema.clicked.connect(self.add_cinema)
        self.addHall.clicked.connect(self.add_hall)
        self.setChairs.clicked.connect(self.set_chairs)
        self.createSession.clicked.connect(self.create_session)
        self.sellTicket.clicked.connect(self.sell_ticket)
        self.closestSession.clicked.connect(self.closest_session)
        self.hallPlan.clicked.connect(self.hall_plan)
        self.seatsInRow.clicked.connect(self.seats_in_row)
        self.sessionsSchedule.clicked.connect(self.sessions_schedule)
        self.cinemaGraph.clicked.connect(self.cinema_graph)
        self.adBooklet.clicked.connect(self.ad_booklet)
        self.visitorsFeedback.clicked.connect(self.visitors_feedback)

    def add_cinema(self):
        self.resultLabel.setText('')
        self.form = AddCinema(self)
        self.form.show()

    def add_hall(self):
        self.resultLabel.setText('')
        if len(self.cinemas) == 0:
            self.resultLabel.setText('Сначала добавьте кинотеатр')
        else:
            self.form = AddHall(self)
            self.form.show()

    def set_chairs(self):
        self.resultLabel.setText('')
        if len(self.cinemas) == 0:
            self.resultLabel.setText('Сначала добавьте кинотеатр')
        elif not any([len(self.cinemas[c]) for c in self.cinemas]):
            self.resultLabel.setText('Сначала добавьте зал в кинотеатр')
        else:
            self.form = SetChairs(self)
            self.form.show()

    def create_session(self):
        self.resultLabel.setText('')
        if len(self.cinemas) == 0:
            self.resultLabel.setText('Сначала добавьте кинотеатр')
        elif not any([len(self.cinemas[c]) for c in self.cinemas]):
            self.resultLabel.setText('Сначала добавьте зал в кинотеатр')
        else:
            self.form = CreateSession(self)
            self.form.show()

    def sell_ticket(self):
        self.resultLabel.setText('')
        if len(self.sessions) == 0:
            self.resultLabel.setText('Сначала создайте сеанс')
        else:
            self.form = SellTicket(self)
            self.form.show()

    def closest_session(self):
        self.resultLabel.setText('')
        if len(self.sessions) == 0:
            self.resultLabel.setText('Сначала создайте сеанс')
            return
        films = set()
        for s in self.sessions:
            films.add(s['film'])

        film, ok_pressed = QInputDialog.getItem(self, 'Введите название',
                                                "Сеанс какого фильма показываем?",
                                                films, 0,
                                                False)
        if not ok_pressed:
            return

        closest = None
        for s in self.sessions:
            if s['film'] == film and s['date'] >= datetime.now().date():
                if not closest:
                    closest = s
                elif s['date'] < closest['date']:
                    closest = s
                elif s['date'] == closest['date'] and s['start'] < closest['start']:
                    closest = s
        start = f'{str(closest["start"].hour()).rjust(2, "0")}:{str(closest["start"].minute()).rjust(2, "0")}'

        if closest:
            self.resultLabel.setText(
                f'Ближайший сеанс фильма "{closest["film"]}"\nпройдёт в зале №{closest["hall_num"]}\n'
                f'кинотеатра "{closest["cinema"]}" {closest["date"].toString("dd.MM.yyyy")} в {start}.')
        else:
            self.resultLabel.setText('Нет сеанса на этот фильм')

    def hall_plan(self):
        self.resultLabel.setText('')
        if len(self.sessions) == 0:
            self.resultLabel.setText('Сначала создайте сеанс')
            return

        sessions_texts = []
        for s in self.sessions:
            sessions_texts.append(session_in_text(s).replace('\n', ' '))

        session, ok_pressed = QInputDialog.getItem(self, 'Выберете сеанс',
                                                   "Выберете сеанс?",
                                                   sessions_texts, 0,
                                                   False)

        session = self.sessions[sessions_texts.index(session)]
        self.form = HallPlan(session['hall'])
        self.form.show()

    def seats_in_row(self):
        self.resultLabel.setText('')
        if len(self.cinemas) == 0:
            self.resultLabel.setText("Сначала добавьте кинотеатр")
        else:
            self.form = SeatsInRow(self)
            self.form.show()

    def sessions_schedule(self):
        self.resultLabel.setText('')
        workbook = xlsxwriter.Workbook('Расписание.xlsx')
        worksheet = workbook.add_worksheet()

        worksheet.write(0, 0, 'Фильм')
        worksheet.write(0, 1, 'Кинотеатр')
        worksheet.write(0, 2, 'Зал')
        worksheet.write(0, 3, 'Дата проведения')
        worksheet.write(0, 4, 'Время начала')
        worksheet.write(0, 5, 'Длительность')

        sorted_sessions = sorted(self.sessions, key=lambda x: x['start'])
        sorted_sessions.sort(key=lambda x: x['date'])

        for row, s in enumerate(sorted_sessions):
            if datetime.now().date() >= s['date'] >= datetime.now().date() - timedelta(days=30):
                worksheet.write(row + 1, 0, s['film'])
                worksheet.write(row + 1, 1, s['cinema'])
                worksheet.write(row + 1, 2, s['hall_num'])
                worksheet.write(row + 1, 3, s['date'].toString("dd.MM.yyyy"))
                start = f'{str(s["start"].hour()).rjust(2, "0")}:{str(s["start"].minute()).rjust(2, "0")}'
                worksheet.write(row + 1, 4, start)
                duration = f'{str(s["duration"].hour()).rjust(2, "0")}:{str(s["duration"].minute()).rjust(2, "0")}'
                worksheet.write(row + 1, 5, duration)

        workbook.close()
        self.resultLabel.setText('Расписание за прошлый месяц здесь: Расписание.xlsx')

    def cinema_graph(self):
        self.resultLabel.setText('')
        cinema, ok_pressed = QInputDialog.getItem(self, 'Выберете кинотеатр',
                                                  "Выберете кинотеатр?",
                                                  list(self.cinemas.keys()), 0,
                                                  False)
        if not ok_pressed:
            return

        workbook = xlsxwriter.Workbook(f'График {cinema}.xlsx')
        worksheet = workbook.add_worksheet()

        # Данные
        data = {}

        for i in range(24):
            data[i] = 0

        for s in self.sessions:
            if s['cinema'] == cinema:
                taked = 0
                for row in s['hall']:
                    taked += sum(row)
                data[s['start'].hour()] += taked

        worksheet.write_column('A1', [f'{str(t).rjust(2, "0")}:00' if t % 3 == 0 else '' for t in data.keys()])
        worksheet.write_column('B1', data.values())

        # Тип диаграммы
        chart = workbook.add_chart({'type': 'line'})

        chart.set_x_axis({
            'name': 'Время',
            'name_font': {
                'name': 'Courier New',
                'color': '#92D050'
            },
            'num_font': {
                'name': 'Arial',
                'color': '#00B0F0',
            },
        })

        chart.set_y_axis({
            'name': 'Занято мест',
            'name_font': {
                'name': 'Courier New',
                'color': '#92D050'
            },
            'num_font': {
                'name': 'Arial',
                'color': '#00B0F0',
            },
        })

        chart.set_size({'width': 900, 'height': 400})

        # Строим по нашим данным
        chart.add_series({'values': f'=Sheet1!B1:B{len(data)}', 'categories': f'=Sheet1!A1:A{len(data)}'})

        worksheet.insert_chart('D1', chart)
        workbook.close()

        self.resultLabel.setText(f'График здесь: График {cinema}.xlsx')

    def ad_booklet(self):
        self.resultLabel.setText('')
        self.form = AdBooklet(self)
        self.form.show()

    def visitors_feedback(self):
        self.resultLabel.setText('')
        first_names = ['Герман', 'Руслан', 'Николай',
                       'Леонид', 'Роман', 'Альберт', 'Анатолий',
                       'Эдуард', 'Юрий', 'Владислав']
        last_names = ['Столичный', 'Амфилохов', 'Ефиманов',
                      'Брыластов', 'Дзасохов', 'Чапкин', 'Михелев',
                      'Пелымсих', 'Чичканов', 'Кузнецов']

        feedbacks = [
            "Кинотеатр просто волшебный! Уютная атмосфера, комфортные кресла и прекрасное качество звука"
            " и изображения. Всегда радуют широкий выбор фильмов для любого вкуса.",
            "Часто посещаю этот кинотеатр и всегда остаюсь довольным. Вежливый персонал,"
            " быстрая обслуживание и чистота в зале - все на высшем уровне.",
            "Отличный выбор фильмов для детей! Детская комната и мягкие кресла позволяют"
            " нам расслабиться и насладиться фильмом, не беспокоясь о наших малышах.",
            "Превосходное звуковое и видеооборудование! Кина технологии настолько качественные,"
            " что кажется, будто попал внутрь фильма. Очень рекомендую!",
            "Один из самых удобных кинотеатров, где я был! Просторные имягкие кресла, несколько"
            " вариантов закусок и напитков - все, что нужно, чтобы максимально насладиться просмотром фильма.",
            "Фантастическая акустика! Звук охватывает весь зал и создает эффект полного погружения"
            " в кино. Настоящая находка для киноманов.",
            "Кинотеатр с душой! Здесь всегда проходят различные тематические мероприятия,"
            " конкурсы и прочие развлечения, добавляющие праздничную атмосферу.",
            "Особое спасибо за удобный онлайн-бронирование билетов! Никогда не приходится"
            " тратить время на ожидание в очередях, все быстро и удобно.",
            "Семейный кинотеатр, который радует и детей, и взрослых. Великолепный выбор фильмов"
            " для разных возрастных категорий и приятная атмосфера для приятного времяпровождения всей семьей.",
            "Я уже не представляю свою жизнь без этого кинотеатра. Здесь всегда превосходный сервис,"
            " отличная киноафиша и уютная атмосфера. Лучшее место для отдыха с друзьями в выходные!"]

        prs = Presentation()

        # получаем схему расположения элементов для заголовочного слайда
        title_slide_layout = prs.slide_layouts[0]
        # создаем заголовочный слайд
        slide = prs.slides.add_slide(title_slide_layout)
        # создаем у слайда заголовок и текст
        title = slide.shapes.title
        title.text = "Отзывы посетителей"
        subtitle = slide.placeholders[1]
        subtitle.text = "Отзывы генерируюся случайно"

        avatars = list(range(1, 11))

        for i in range(5):
            feedback = choice(feedbacks)
            feedbacks.remove(feedback)

            full_name = [choice(first_names), choice(last_names)]
            first_names.remove(full_name[0])
            last_names.remove(full_name[1])
            full_name = ' '.join(full_name)

            slide = prs.slides.add_slide(prs.slide_layouts[8])
            slide.shapes.title.text = full_name
            avatar = choice(avatars)
            avatars.remove(avatar)
            slide.placeholders[1].insert_picture(f'avatars/avatar{avatar}.png')
            slide.placeholders[2].text = feedback

        prs.save('Отзывы.pptx')
        self.resultLabel.setText("Отзывы здесь: Отзывы.pptx")


class AddCinema(QMainWindow):
    def __init__(self, parent):
        super().__init__()
        uic.loadUi("addCinema.ui", self)
        self.parent = parent
        self.addBtn.clicked.connect(self.add)

    def add(self):
        cinema = self.titleInput.text()
        if cinema in self.parent.cinemas:
            self.statusBar().showMessage('Такой кинотеатр уже есть', 5000)
        else:
            self.parent.cinemas[cinema] = []
            self.close()
            self.parent.resultLabel.setText('Кинотеатр добавлен')


def session_in_text(s):
    start = f'{str(s["start"].hour()).rjust(2, "0")}:{str(s["start"].minute()).rjust(2, "0")}'
    duration = f'{str(s["duration"].hour()).rjust(2, "0")}:{str(s["duration"].minute()).rjust(2, "0")}'
    return f'Фильм "{s["film"]}",\nдата проведения {s["date"].toString("dd.MM.yyyy")} в {start},\n' \
           f'длительност {duration}, зал №{s["hall_num"]} кинотеатра "{s["cinema"]}"'


def except_hook(cls, exception, traceback):
    sys.excepthook(cls, exception, traceback)


if __name__ == '__main__':
    app = QApplication(sys.argv)
    mcs = MyCinemasSystem()
    mcs.show()
    sys.excepthook = except_hook
    sys.exit(app.exec())
